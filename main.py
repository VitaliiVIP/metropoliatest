import os
import json
import random
from uuid import uuid4
from flask import Flask, render_template, request, session
from openai import OpenAI
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation

load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "change-this-secret")
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

quiz_states = {}

pdf_text_global = None
chunks_global = []
chunk_index_global = 0
questions = 0
right_answers = 0
questions_amount = 20
final_text = ""

def get_quiz_state():
    # Ensure each browser has its own session_id cookie
    sid = session.get("session_id")
    if not sid:
        sid = str(uuid4())
        session["session_id"] = sid

    # Create state if not exists
    if sid not in quiz_states:
        quiz_states[sid] = {
            "pdf_text": None,
            "chunks": [],
            "chunk_index": 0,
            "questions": 0,
            "right_answers": 0,
        }
    return quiz_states[sid]

#prompt
def generate_question(text_chunk):
    prompt = f"""
    You are an AI that generates quiz questions from study material.

    Based ONLY on the following text:

    \"\"\"{text_chunk}\"\"\"

    Create ONE multiple-choice question with EXACTLY:
    - 1 correct answer
    - 1 wrong answer (plausible but incorrect)

    IMPORTANT RULES:
    - Answer ONLY with valid JSON.
    - No explanations.
    - No extra text.
    - No backticks.
    - No markdown.

    JSON format:
    {{
      "question": "...",
      "correct_answer": "...",
      "wrong_answer": "..."
    }}
    """

    response = client.responses.create(
        model="gpt-4.1-mini",
        input=prompt,
    )

    raw_text = getattr(response, "output_text", None)

    if not raw_text:
        try:
            part = response.output[0].content[0].text
            # in some SDK versions .text is an object with .value
            raw_text = getattr(part, "value", part)
        except Exception:
            raw_text = ""

    print("RAW LLM OUTPUT:", repr(raw_text))  # debug

    # Safely parse JSON (no crash if it's bad)
    try:
        data = json.loads(raw_text)
        answers = shuffle_answers(data["correct_answer"], data["wrong_answer"])
        data["answers"] = answers
    except Exception as e:
        print("JSON parse failed:", e)
        # Fallback so the app continues working
        data["answers"] = shuffle_answers(
            data["correct_answer"],
            data["wrong_answer"]
        )
    return data

def shuffle_answers(correct, wrong):
    answers = [
        {"text": correct, "is_correct": True},
        {"text": wrong, "is_correct": False}
    ]
    random.shuffle(answers)
    return answers

#pdftotext
def extract_pdf_text(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_pptx_text(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def split_into_chunks(text, max_chars=1350):
    chunks = []
    current = ""
    for line in text.split("\n"):
        if len(current) + len(line) > max_chars:
            chunks.append(current)
            current = ""
        current += line + "\n"
    if current:
        chunks.append(current)
    return chunks

@app.route("/", methods=["GET", "POST"])
def index():
    state = get_quiz_state()

    if request.method == "POST":
        file = request.files.get("pdf")

        # === New quiz start ===
        if file and file.filename != "":
            filename = file.filename.lower()

            if filename.endswith(".pdf"):
                text = extract_pdf_text(file)
            elif filename.endswith(".pptx"):
                text = extract_pptx_text(file)
            else:
                return "Unsupported file type (use .pdf or .pptx)"

            # reset THIS USER'S state
            state["pdf_text"] = text
            state["chunks"] = split_into_chunks(text)
            state["chunk_index"] = 0
            state["questions"] = 0
            state["right_answers"] = 0

            chunk = state["chunks"][state["chunk_index"] % len(state["chunks"])]
            state["chunk_index"] += 1
            ai_output = generate_question(chunk)

            return render_template(
                "quiz.html",
                data=ai_output,
                result=None,
                questions=state["questions"],
                right_answers=state["right_answers"],
                questions_amount=questions_amount,
            )

        # === Answer to a question ===
        action = request.form.get("action")
        if action == "answer":
            if not state["pdf_text"]:
                return "No file uploaded for this session"

            correct = request.form.get("correct")
            user_answer = request.form.get("answer")

            state["questions"] += 1
            if user_answer == correct:
                state["right_answers"] += 1
                result = "Yes!"
            else:
                result = "No :("

            q = state["questions"]
            r = state["right_answers"]

            if q >= questions_amount:
                score = r / q
                if score == 1:
                    final_text = "Excellent! You are fully prepared with this topic."
                elif score >= 0.8:
                    final_text = "Good! You answered almost all questions!"
                elif score >= 0.55:
                    final_text = "Not bad! Keep going."
                else:
                    final_text = "I would recommend you to spend more time for this topic and try again."

                return render_template(
                    "score.html",
                    questions=q,
                    right_answers=r,
                    questions_amount=questions_amount,
                    final_text=final_text,
                )

            # next question
            if not state["chunks"]:
                state["chunks"] = split_into_chunks(state["pdf_text"])

            chunk = state["chunks"][state["chunk_index"] % len(state["chunks"])]
            state["chunk_index"] += 1
            ai_output = generate_question(chunk)

            return render_template(
                "quiz.html",
                data=ai_output,
                result=result,
                questions=q,
                right_answers=r,
                questions_amount=questions_amount,
            )

        return "Invalid request"

    # GET request
    return render_template("index.html")

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=80, debug=True)
